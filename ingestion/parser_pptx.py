"""
parser_pptx.py — Extract content from PowerPoint files using python-pptx.

Returns a list of slide dicts, each containing:
  - slide_num: int
  - text_blocks: list of text strings
  - image_blocks: list of image bytes (for OCR)
  - has_tables: bool
  - table_data: list of tables (each table is list of rows, each row list of cell strings)
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def parse_pptx(filepath: str) -> list[dict]:
    """
    Parse a PPTX file and return structured blocks per slide.
    """
    prs = Presentation(filepath)
    slides = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_blocks = []
        image_blocks = []
        table_data = []
        has_tables = False

        for shape in slide.shapes:
            # --- Text ---
            if shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        paragraphs.append(text)
                if paragraphs:
                    text_blocks.append({
                        "text": "\n".join(paragraphs),
                        "bbox": (shape.left, shape.top, shape.width, shape.height),
                    })

            # --- Images ---
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_blob = shape.image.blob
                    image_blocks.append({
                        "image_bytes": img_blob,
                        "bbox": (shape.left, shape.top, shape.width, shape.height),
                    })
                except Exception:
                    continue

            # --- Tables ---
            if shape.has_table:
                has_tables = True
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                table_data.append(rows)

        slides.append({
            "page_num": slide_num,
            "text_blocks": text_blocks,
            "image_blocks": image_blocks,
            "has_tables": has_tables,
            "table_data": table_data,
        })

    return slides
